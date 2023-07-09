import json
import os
import re
import time
import asyncio
from pathlib import Path
from datetime import datetime
import openai
import pptx

from db import create_session
from model import Upload

EXIT_SUCCESS = 0
CONTENT = [
    {"role": "system", "content": "Can you explain the slides in basic English, and provide examples if needed!"}
]
ERROR_MESSAGE = "Something is wrong:"
ENGINE_MODEL = "gpt-3.5-turbo"
WRITE_TO_FILE_MODE = 'w'
UPLOADS_DIR = "uploads"
OUTPUTS_DIR = "outputs"


session = create_session()


async def parse_file_to_slides(file_path):
    """
    Parse the PowerPoint presentation file and extract slides.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        list: List of slides, where each slide is a list of text content.
    """
    ppt = await asyncio.to_thread(pptx.Presentation, file_path)
    slides = []

    for slide in ppt.slides:
        slide_content = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_content.append(run.text)

        slides.append(slide_content)

    return slides


async def explain_slides(slides):
    """
    Generate explanations for each slide.

    Args:
        slides (list): List of slides, where each slide is a list of text content.

    Returns:
        list: List of explanations for each slide.
    """
    explanations = []

    for slide in slides:
        explanation = await parse_slide_of_pptx(slide)
        explanations.append(explanation)

    return explanations


async def parse_slide_of_pptx(slide):
    """
    Parse a single slide of a PowerPoint presentation.

    Args:
        slide (list): List of text content in the slide.

    Returns:
        str: Explanation of the slide.
    """
    try:
        response = await request_completion(slide)
        return response
    except Exception as error:
        error_message = f"{ERROR_MESSAGE} Error processing slide: {str(error)}"
        return error_message


async def request_completion(slide):
    """
    Send slide content to OpenAI's Chat API to generate explanations.

    Args:
        slide (list): List of text content in the slide.

    Returns:
        str: Generated explanation for the slide.
    """
    CONTENT.append({"role": "user", "content": " ".join(slide)})
    try:
        response = openai.ChatCompletion.create(
            model=ENGINE_MODEL,
            messages=CONTENT
        )
        content = response["choices"][0].message.content
        cleaned_text = re.sub(r"[\n\r]", "", content).encode("ascii", "ignore").decode("utf-8")
        return cleaned_text.strip()
    except Exception as error:
        print(f"Error during completion request: {str(error)}")
        return None


def modify_file(explanations, presentation_path):
    """
    Modify the file path and save explanations as JSON.

    Args:
        explanations (list): List of explanations for each slide.
        presentation_path (str): Path to the original presentation file.
    """
    try:
        presentation_name = os.path.basename(presentation_path)
        presentation_name = os.path.splitext(presentation_name)[0]
        output_file = os.path.join(OUTPUTS_DIR, f"{presentation_name}.json")

        # Create a dictionary with slide numbers as keys and explanations as values
        slide_explanations = {f"slide{slide_num}": explanation for slide_num, explanation in enumerate(explanations, start=1)}
        # Save the slide explanations as JSON
        with open(output_file, WRITE_TO_FILE_MODE) as file:
            json.dump(slide_explanations, file, indent=4)
            print(f"Explanations saved to {output_file}")
    except FileNotFoundError as error:
        print(f"{ERROR_MESSAGE} Error saving explanations: {str(error)}")


async def process_file(upload_id):
    """
    Process a single file by parsing slides, generating explanations, and saving the results.

    Args:
        file_path (str): Path to the PowerPoint file.
    """

    start_time = time.time()

    upload = session.get(Upload, upload_id)

    if upload is None:
        print(f"Upload with ID '{upload_id}' not found")
        return

    presentation_path = f"{UPLOADS_DIR}/{upload.uid}.pptx"
    print(presentation_path)

    # Parse the file to extract slides
    slides = await parse_file_to_slides(presentation_path)
    # Generate explanations for each slide
    explanations = await explain_slides(slides)
    # Save the explanations to a file
    modify_file(explanations, presentation_path)
    end_time = time.time()

    upload.status = 'completed'
    upload.finish_time = datetime.now()

    session.commit()

    execution_time = end_time - start_time
    minutes, seconds = divmod(execution_time, 60)
    print(f"Execution time: {minutes:.0f} minutes {seconds:.2f} seconds")



async def process_files_in_uploads():
    """
    Process files in the uploads directory in an infinite loop.
    """
    # uploads_path = Path(UPLOADS_DIR)
    # processed_files = set()

    while True:

        pending_uploads = session.query(Upload).filter_by(status="pending").all()

        for upload in pending_uploads:
            upload_id = upload.id
            await process_file(upload_id)

        await asyncio.sleep(10)


async def main():
    """
    Main function to initiate file processing and setup necessary directories.

    Returns:
        int: Exit status (0 for success).
    """
    openai.api_key = 'API_KEY'

    # Create the uploads and outputs directories if they don't exist
    Path(UPLOADS_DIR).mkdir(parents=True, exist_ok=True)
    Path(OUTPUTS_DIR).mkdir(parents=True, exist_ok=True)

    await process_files_in_uploads()

    return EXIT_SUCCESS


if __name__ == "__main__":
    asyncio.run(main())
